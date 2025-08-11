# pages/5_Exportar_y_Sesiones.py
# Export & Sessions page (Streamlit multipage)
# IMPORTANT: Code in English. UI labels/messages in Spanish.

from __future__ import annotations

import io
import json
import time
import zipfile
import unicodedata
from dataclasses import asdict
from typing import Dict, List, Tuple

import numpy as np
import pandas as pd
import streamlit as st

# Optional libs (guarded imports)
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.metrics.pairwise import cosine_similarity
except Exception:
    TfidfVectorizer = None
    cosine_similarity = None


# ---------------------------------------
# Helpers
# ---------------------------------------
def normalize_text(s: str) -> str:
    s = str(s)
    s = s.lower()
    s = "".join(c for c in unicodedata.normalize("NFKD", s) if not unicodedata.combining(c))
    return " ".join(s.split())


def flatten_taxonomy(nodes: List[dict], market: str) -> pd.DataFrame:
    rows = []
    for n in nodes:
        rows.append({
            "market": market,
            "id": n.get("id", ""),
            "parent_id": n.get("parent_id") or "",
            "name": n.get("name", ""),
            "slug": n.get("slug", ""),
            "intent": n.get("intent", ""),
            "ms": int(n.get("ms", 0) or 0),
            "score": float(n.get("score", 0.0) or 0.0)
        })
    return pd.DataFrame(rows)


def flatten_filters(nodes: List[dict], market: str) -> pd.DataFrame:
    rows = []
    for n in nodes:
        for f in (n.get("filters") or []):
            # handle dict or dataclass-like
            name = f.get("name") if isinstance(f, dict) else getattr(f, "name", "")
            ftype = f.get("type") if isinstance(f, dict) else getattr(f, "type", "in")
            values = f.get("values") if isinstance(f, dict) else getattr(f, "values", [])
            order = f.get("order") if isinstance(f, dict) else getattr(f, "order", 0)
            visibility = f.get("visibility") if isinstance(f, dict) else getattr(f, "visibility", "core")
            rows.append({
                "market": market,
                "category_id": n.get("id", ""),
                "name": name,
                "type": ftype,
                "values": "|".join(values or []),
                "order": int(order or 0),
                "visibility": visibility,
            })
    return pd.DataFrame(rows) if rows else pd.DataFrame(columns=["market","category_id","name","type","values","order","visibility"])


def compute_kw_to_node(df_kw: pd.DataFrame, nodes: List[dict]) -> pd.DataFrame:
    """
    Map keywords to best child node using TF-IDF cosine. Falls back to token overlap.
    Expects df_kw with columns ['keyword','kw_norm'].
    """
    # Identify child nodes
    root = next((n for n in nodes if n.get("parent_id") is None), None)
    children = [n for n in nodes if n.get("parent_id") == (root.get("id") if root else None)]
    if df_kw.empty or not children:
        return pd.DataFrame(columns=["keyword", "node_id", "similarity", "intent"])

    child_names = [c["name"] for c in children]

    if TfidfVectorizer is not None and cosine_similarity is not None:
        texts = df_kw["kw_norm"].tolist() + child_names
        vec = TfidfVectorizer(ngram_range=(1, 2), min_df=1)
        X = vec.fit_transform(texts)
        A = X[:len(df_kw)]
        B = X[len(df_kw):]
        sims = cosine_similarity(A, B)
        idx = np.argmax(sims, axis=1)
        best = sims[np.arange(sims.shape[0]), idx]
    else:
        # Fallback: token overlap
        def overlap(a: str, b: str) -> float:
            ta, tb = set(normalize_text(a).split()), set(normalize_text(b).split())
            if not ta and not tb: return 0.0
            return len(ta & tb) / max(1, len(ta))
        idx, best = [], []
        for q in df_kw["kw_norm"].tolist():
            sims = [overlap(q, name) for name in child_names]
            j = int(np.argmax(sims)) if sims else 0
            idx.append(j)
            best.append(sims[j] if sims else 0.0)

    # Simple intent heuristic
    def intent_of(q: str) -> str:
        qn = normalize_text(q)
        if any(k in qn for k in ["precio","oferta","comprar","barat","mejor","calidad precio",
                                  "preço","oferta",
                                  "acheter","prix","promo",
                                  "prezzo","comprare","sconto"]):
            return "transactional"
        if any(k in qn for k in ["como","cómo","que es","qué es","cual","cuál",
                                  "difference","diferencia","come","cosa è","o que é"]):
            return "informational"
        return "mixed"

    rows = []
    for (kw, qn), j, s in zip(df_kw[["keyword","kw_norm"]].itertuples(index=False), idx, best):
        node = children[int(j)]
        rows.append({"keyword": kw, "node_id": node["id"], "similarity": round(float(s), 3), "intent": intent_of(kw)})
    return pd.DataFrame(rows)


def export_master_json_per_market(nodes_by_market: Dict[str, List[dict]]) -> str:
    """
    Master JSON as a dict: { "es": [nodes...], "pt": [...], ... }
    """
    safe = {mk: nodes for mk, nodes in nodes_by_market.items()}
    return json.dumps(safe, ensure_ascii=False, indent=2)


def build_markdown_summary(nodes_by_market: Dict[str, List[dict]]) -> bytes:
    md = io.StringIO()
    md.write("# Resumen ejecutivo — Taxonomía propuesta\n\n")
    for mk, nodes in nodes_by_market.items():
        root = next((n for n in nodes if n.get("parent_id") is None), None)
        children = [n for n in nodes if n.get("parent_id") == (root.get("id") if root else None)] if root else []
        md.write(f"## Mercado {mk.upper()}\n```\n")
        if root:
            md.write(root["name"] + "\n")
            for c in children:
                md.write(f"└── {c['name']}\n")
        else:
            md.write("(Árbol no disponible)\n")
        md.write("```\n\n")
    return md.getvalue().encode("utf-8")


def build_pptx_summary(nodes_by_market: Dict[str, List[dict]]) -> bytes:
    if not PPTX_AVAILABLE:
        return b""
    prs = Presentation()
    t = prs.slides.add_slide(prs.slide_layouts[0])
    t.shapes.title.text = "Resumen ejecutivo — Taxonomía"
    t.placeholders[1].text = "Arquitectura por mercado"

    for mk, nodes in nodes_by_market.items():
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = f"Arquitectura — {mk.upper()}"
        tf = s.placeholders[1].text_frame
        root = next((n for n in nodes if n.get("parent_id") is None), None)
        children = [n for n in nodes if n.get("parent_id") == (root.get("id") if root else None)] if root else []
        if root:
            tf.text = root["name"]
            for c in children:
                p = tf.add_paragraph()
                p.text = f"└── {c['name']}  ({c.get('ms',0)} ms, score {c.get('score',0)})"
        else:
            tf.text = "(Árbol no disponible)"

    mem = io.BytesIO()
    prs.save(mem)
    mem.seek(0)
    return mem.read()


def package_zip(files: List[Tuple[str, bytes]]) -> bytes:
    mem = io.BytesIO()
    with zipfile.ZipFile(mem, "w", compression=zipfile.ZIP_DEFLATED) as z:
        for fname, data in files:
            z.writestr(fname, data)
    mem.seek(0)
    return mem.read()


# ---------------------------------------
# Preconditions
# ---------------------------------------
st.header("6) Exportar & Sesiones")
st.caption("Descarga JSON/CSVs/PPTX/Markdown y guarda/carga sesiones de trabajo.")

if "taxonomy_nodes" not in st.session_state or not st.session_state["taxonomy_nodes"]:
    st.error("No hay una jerarquía propuesta en sesión. Ve a **3) Clustering y Taxonomía** primero.")
    st.stop()

params = st.session_state.get("params", {})
markets_available = sorted(st.session_state["taxonomy_nodes"].keys())
markets_selected = st.multiselect("Mercados a exportar", markets_available, default=markets_available)

# Optional sources for extras
gkp_by_market: Dict[str, pd.DataFrame] = st.session_state.get("gkp_data", {})
clusters_preview: Dict[str, pd.DataFrame] = st.session_state.get("clusters_preview", {})

st.divider()
st.subheader("Export individual por mercado")

export_rows = []
zip_payload: List[Tuple[str, bytes]] = []

for mk in markets_selected:
    nodes = st.session_state["taxonomy_nodes"].get(mk, [])
    if not nodes:
        st.warning(f"[{mk.upper()}] No hay nodos para exportar.")
        continue

    st.markdown(f"**{mk.upper()}**")

    # JSON (per market)
    json_bytes = json.dumps(nodes, ensure_ascii=False, indent=2).encode("utf-8")
    st.download_button(f"Descargar master_{mk}.json", data=json_bytes, file_name=f"master_{mk}.json", mime="application/json")
    zip_payload.append((f"{mk}/master_{mk}.json", json_bytes))

    # taxonomy.csv (per market)
    tax_df = flatten_taxonomy(nodes, mk)
    tax_bytes = tax_df.to_csv(index=False).encode("utf-8")
    st.download_button(f"Descargar taxonomy_{mk}.csv", data=tax_bytes, file_name=f"taxonomy_{mk}.csv", mime="text/csv")
    zip_payload.append((f"{mk}/taxonomy_{mk}.csv", tax_bytes))

    # filters.csv (per market)
    filt_df = flatten_filters(nodes, mk)
    filt_bytes = filt_df.to_csv(index=False).encode("utf-8")
    st.download_button(f"Descargar filters_{mk}.csv", data=filt_bytes, file_name=f"filters_{mk}.csv", mime="text/csv")
    zip_payload.append((f"{mk}/filters_{mk}.csv", filt_bytes))

    # keywords_to_node.csv (per market) – use clustered/normalized keywords if available
    df_kw = clusters_preview.get(mk, gkp_by_market.get(mk, pd.DataFrame())).copy()
    if not df_kw.empty:
        if "kw_norm" not in df_kw.columns:
            df_kw["kw_norm"] = df_kw["keyword"].astype(str).map(normalize_text)
        kw_map = compute_kw_to_node(df_kw[["keyword","kw_norm"]], nodes)
    else:
        kw_map = pd.DataFrame(columns=["keyword","node_id","similarity","intent"])
    kw_bytes = kw_map.to_csv(index=False).encode("utf-8")
    st.download_button(f"Descargar keywords_to_node_{mk}.csv", data=kw_bytes, file_name=f"keywords_to_node_{mk}.csv", mime="text/csv")
    zip_payload.append((f"{mk}/keywords_to_node_{mk}.csv", kw_bytes))

st.divider()
st.subheader("Export agregado (todos los mercados seleccionados)")

# Master JSON (all markets)
master_json = export_master_json_per_market({mk: st.session_state["taxonomy_nodes"][mk] for mk in markets_selected})
st.download_button("Descargar master_all.json", data=master_json.encode("utf-8"), file_name="master_all.json", mime="application/json")
zip_payload.append(("master_all.json", master_json.encode("utf-8")))

# Markdown summary (all)
md_bytes = build_markdown_summary({mk: st.session_state["taxonomy_nodes"][mk] for mk in markets_selected})
st.download_button("Descargar resumen (Markdown)", data=md_bytes, file_name="resumen_taxonomia.md", mime="text/markdown")
zip_payload.append(("resumen_taxonomia.md", md_bytes))

# PPTX (optional)
if PPTX_AVAILABLE:
    pptx_bytes = build_pptx_summary({mk: st.session_state["taxonomy_nodes"][mk] for mk in markets_selected})
    st.download_button("Descargar resumen ejecutivo (PPTX)", data=pptx_bytes, file_name="resumen_taxonomia.pptx",
                       mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    zip_payload.append(("resumen_taxonomia.pptx", pptx_bytes))
else:
    st.info("PPTX no disponible: instala `python-pptx` para habilitar esta exportación.")

# ZIP all
ts = int(time.time())
zip_bytes = package_zip(zip_payload) if zip_payload else b""
st.download_button("Descargar paquete completo (ZIP)", data=zip_bytes, file_name=f"exports_{ts}.zip", mime="application/zip")

st.divider()
st.subheader("Sesiones")

# Save session JSON (compact but useful)
session_snapshot = {
    "params": params,
    "markets": markets_selected,
    "taxonomy_nodes": {mk: st.session_state["taxonomy_nodes"][mk] for mk in markets_selected},
    "coverage_by_market": st.session_state.get("coverage_by_market", pd.DataFrame()).to_dict(orient="list")
        if isinstance(st.session_state.get("coverage_by_market"), pd.DataFrame) else st.session_state.get("coverage_by_market", {}),
    "timestamp": ts,
}
sess_bytes = json.dumps(session_snapshot, ensure_ascii=False, indent=2).encode("utf-8")
st.download_button("Guardar sesión (.json)", data=sess_bytes, file_name=f"session_{ts}.json", mime="application/json")

with st.expander("Cargar sesión desde archivo (.json)"):
    uploaded = st.file_uploader("Selecciona la sesión", type=["json"], key="session_loader_5")
    if uploaded is not None:
        try:
            data = json.load(uploaded)
            # Restore minimal required keys
            if "taxonomy_nodes" in data:
                for mk, nodes in data["taxonomy_nodes"].items():
                    st.session_state.setdefault("taxonomy_nodes", {})
                    st.session_state["taxonomy_nodes"][mk] = nodes
            if "params" in data:
                st.session_state["params"] = data["params"]
            st.success("Sesión cargada. Vuelve a las páginas anteriores para continuar la edición o re-exportar.")
        except Exception as e:
            st.error(f"No se pudo cargar la sesión: {e}")

st.caption("Sugerencia: versiona los `master_*.json` y `taxonomy_*.csv` en Git para auditar cambios entre iteraciones.")
