# streamlit_app.py
# Multilingual e-commerce taxonomy & filters generator (Technology vertical)
# IMPORTANT: All code in English. UI labels/messages in Spanish as requested.

import os
import io
import json
import uuid
import math
import time
import zipfile
import unicodedata
from dataclasses import dataclass, field, asdict
from typing import List, Dict, Any, Optional, Tuple

import pandas as pd
import numpy as np
import streamlit as st

# Optional libs (guarded imports)
try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.cluster import AgglomerativeClustering
    from sklearn.metrics.pairwise import cosine_similarity
except Exception:
    TfidfVectorizer = None
    AgglomerativeClustering = None
    cosine_similarity = None

try:
    import plotly.express as px
except Exception:
    px = None

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

# -----------------------------
# App config
# -----------------------------
st.set_page_config(page_title="Generador de Taxonomías SEO (Tech)", layout="wide")


# -----------------------------
# Data models (schemas)
# -----------------------------
@dataclass
class SEOData:
    title: str
    h1: str
    meta: str
    faq: List[str] = field(default_factory=list)


@dataclass
class FilterDef:
    name: str
    type: str  # "in", "range", "boolean"
    values: List[str] = field(default_factory=list)
    order: int = 0
    visibility: str = "core"  # "core" | "advanced" | "hidden_mobile"


@dataclass
class Node:
    id: str
    parent_id: Optional[str]
    name: str
    slug: str
    intent: str  # "informational" | "mixed" | "transactional"
    ms: int  # monthly searches aggregated for node
    score: float
    recommended_PLPs: List[str] = field(default_factory=list)
    filters: List[FilterDef] = field(default_factory=list)
    seo: SEOData = field(default_factory=lambda: SEOData("", "", "", []))


# -----------------------------
# Helpers (locale & text)
# -----------------------------
SUPPORTED_MARKETS = ["es", "pt", "fr", "it"]

LOCALIZED_ROOT_DEFAULT = {
    "es": "Tecnología",
    "pt": "Tecnologia",
    "fr": "Technologie",
    "it": "Tecnologia",
}

# Simple locale-aware strings
I18N = {
    "download_json": {
        "es": "Descargar master_{mk}.json"
    },
    "download_tax": {
        "es": "Descargar taxonomy_{mk}.csv"
    },
    "download_filters": {
        "es": "Descargar filters_{mk}.csv"
    },
    "download_kw_map": {
        "es": "Descargar keywords_to_node_{mk}.csv"
    },
    "download_md": {
        "es": "Descargar resumen (Markdown)"
    },
    "download_pptx": {
        "es": "Descargar resumen ejecutivo (PPTX)"
    },
}


def normalize_text(s: str) -> str:
    s = str(s)
    s = s.lower()
    s = "".join(c for c in unicodedata.normalize("NFKD", s) if not unicodedata.combining(c))
    return " ".join(s.split())


def make_slug(name: str) -> str:
    s = normalize_text(name).replace(" ", "-")
    s = "".join(ch for ch in s if ch.isalnum() or ch in "-/")
    return s.strip("-")


def gen_id() -> str:
    return str(uuid.uuid4())


# -----------------------------
# Presets
# -----------------------------
PRESETS = {
    "Calidad (por defecto)": {
        "clustering": "embeddings+hdbscan",  # TODO (stubbed as TF-IDF fallback)
        "use_llm_for_intent_arbiter": True,  # stubbed switch
        "use_llm_for_metadata": True,        # stubbed switch
        "temperature": 0.0,
        "cannibal_jaccard": 0.6,
        "cannibal_cosine": 0.85,
        "max_levels": 3,
        "max_items_per_level": 12,
        "ms_quantile_es": 0.60,
        "ms_quantile_others": 0.50,
    },
    "Rápido": {
        "clustering": "tfidf+agglomerative",
        "use_llm_for_intent_arbiter": False,
        "use_llm_for_metadata": False,
        "temperature": 0.0,
        "cannibal_jaccard": 0.6,
        "cannibal_cosine": 0.85,
        "max_levels": 3,
        "max_items_per_level": 12,
        "ms_quantile_es": 0.60,
        "ms_quantile_others": 0.50,
    },
    "Económico": {
        "clustering": "tfidf+agglomerative",
        "use_llm_for_intent_arbiter": False,
        "use_llm_for_metadata": False,
        "temperature": 0.0,
        "cannibal_jaccard": 0.6,
        "cannibal_cosine": 0.85,
        "max_levels": 3,
        "max_items_per_level": 12,
        "ms_quantile_es": 0.60,
        "ms_quantile_others": 0.50,
    },
    "Avanzado": {
        "clustering": "custom",
        "use_llm_for_intent_arbiter": True,
        "use_llm_for_metadata": True,
        "temperature": 0.0,
        "cannibal_jaccard": 0.6,
        "cannibal_cosine": 0.85,
        "max_levels": 3,
        "max_items_per_level": 12,
        "ms_quantile_es": 0.60,
        "ms_quantile_others": 0.50,
    },
}


# -----------------------------
# File loaders
# -----------------------------
@st.cache_data(show_spinner=False)
def read_gkp_any(file) -> pd.DataFrame:
    """
    Robust GKP parser (CSV/Excel, UTF-16/tab or UTF-8/comma).
    Returns DataFrame with ['keyword', 'avg_monthly_searches'].
    """
    if file is None:
        return pd.DataFrame(columns=["keyword", "avg_monthly_searches"])
    try:
        if file.name.endswith(".xlsx"):
            df = pd.read_excel(file)
        else:
            # try utf-16 + tab
            file.seek(0)
            try:
                df = pd.read_csv(file, encoding="utf-16", sep="\t")
            except Exception:
                file.seek(0)
                df = pd.read_csv(file, encoding="utf-8", sep=",")
    except Exception:
        file.seek(0)
        df = pd.read_csv(file, engine="python")

    df.columns = [str(c).strip().lower() for c in df.columns]
    kw_col = next(
        (c for c in df.columns if any(k in c for k in ["keyword", "palabra", "term", "termo", "terme", "parola", "consulta"])),
        df.columns[0],
    )
    df.rename(columns={kw_col: "keyword"}, inplace=True)
    ms_col = next(
        (c for c in df.columns if ("avg" in c and "search" in c) or ("promedio" in c) or ("moyenne" in c) or ("média" in c) or ("media" in c and "ricer" in c) or ("monthly" in c and "search" in c)),
        None,
    )
    if ms_col:
        df["avg_monthly_searches"] = pd.to_numeric(
            df[ms_col].astype(str).str.replace('"', "", regex=False).str.replace(".", "", regex=False).str.replace(",", "", regex=False),
            errors="coerce",
        ).fillna(0).astype(int)
    else:
        df["avg_monthly_searches"] = 0

    df["keyword"] = df["keyword"].astype(str).str.strip()
    return df[["keyword", "avg_monthly_searches"]]


def read_gsc_zip(upload) -> Tuple[pd.DataFrame, pd.DataFrame]:
    """
    Load GSC export ZIP: returns (queries_df, pages_df).
    """
    if upload is None:
        return pd.DataFrame(), pd.DataFrame()
    with zipfile.ZipFile(upload) as zf:
        qname = next((n for n in zf.namelist() if any(k in n for k in ["Consulta", "Consultas", "query", "queries"])), None)
        pname = next((n for n in zf.namelist() if any(k in n for k in ["Página", "Páginas", "page", "pages", "urls"])), None)
        qdf = pd.read_csv(zf.open(qname)) if qname else pd.DataFrame()
        pdf = pd.read_csv(zf.open(pname)) if pname else pd.DataFrame()

    if not qdf.empty:
        qdf.columns = [c.strip().lower().replace(" ", "_") for c in qdf.columns]
    if not pdf.empty:
        pdf.columns = [c.strip().lower().replace(" ", "_") for c in pdf.columns]
    return qdf, pdf


# -----------------------------
# Core pipeline (POC-friendly, deterministic)
# -----------------------------
def normalize_keywords(df: pd.DataFrame) -> pd.DataFrame:
    d = df.copy()
    d["kw_norm"] = d["keyword"].map(normalize_text)
    return d.drop_duplicates(subset=["kw_norm"])


def cluster_keywords(df: pd.DataFrame, max_items: int) -> pd.DataFrame:
    """
    Lightweight, deterministic clustering with TF-IDF + Agglomerative as fallback.
    Produces cluster_id and cluster_head per keyword.
    """
    d = df.copy()
    if d.empty:
        d["cluster_id"] = []
        d["cluster_head"] = []
        return d

    if TfidfVectorizer is None or AgglomerativeClustering is None or len(d) < 8:
        # Minimal fallback: single cluster
        d["cluster_id"] = 0
        head = d.sort_values("avg_monthly_searches", ascending=False)["keyword"].head(1).tolist()
        d["cluster_head"] = head[0] if head else "General"
        return d

    texts = d["kw_norm"].tolist()
    vec = TfidfVectorizer(ngram_range=(1, 2), min_df=2)
    X = vec.fit_transform(texts)
    n_clusters = max(2, min(max_items, int(math.sqrt(len(texts)))))
    model = AgglomerativeClustering(n_clusters=n_clusters, linkage="ward")
    labels = model.fit_predict(X.toarray())
    d["cluster_id"] = labels

    # head term per cluster by total ms
    heads = {}
    for cid in sorted(set(labels)):
        sub = d[d["cluster_id"] == cid].sort_values("avg_monthly_searches", ascending=False)
        heads[cid] = sub["keyword"].iloc[0] if not sub.empty else f"Cluster {cid}"
    d["cluster_head"] = d["cluster_id"].map(heads)
    return d


def classify_intent(df: pd.DataFrame) -> pd.DataFrame:
    def heuristic(q: str) -> str:
        qn = normalize_text(q)
        if any(k in qn for k in ["precio", "oferta", "comprar", "barat", "mejor", "calidad precio",
                                  "preço", "oferta",
                                  "acheter", "prix", "promo",
                                  "prezzo", "comprare", "sconto"]):
            return "transactional"
        if any(k in qn for k in ["como", "cómo", "que es", "qué es", "cual", "cuál",
                                  "difference", "diferencia", "come", "cosa è", "o que é"]):
            return "informational"
        return "mixed"

    d = df.copy()
    d["intent"] = d["kw_norm"].map(heuristic)
    return d


def propose_hierarchy_for_market(
    df: pd.DataFrame,
    market: str,
    top_level_name: str,
    max_items_per_level: int,
) -> List[Node]:
    """
    Build a 2-level hierarchy:
    - Level 1: user-provided top-level (default 'Tecnología')
    - Level 2: cluster heads (up to max_items_per_level)
    """
    nodes: List[Node] = []
    root = Node(
        id=gen_id(),
        parent_id=None,
        name=top_level_name,
        slug=f"/{make_slug(top_level_name)}",
        intent="mixed",
        ms=int(df["avg_monthly_searches"].sum()),
        score=1.0,
    )
    nodes.append(root)

    # choose up to N cluster heads as subcategories
    candidates = (
        df.groupby(["cluster_id", "cluster_head"], as_index=False)["avg_monthly_searches"]
        .sum()
        .sort_values("avg_monthly_searches", ascending=False)
        .head(max_items_per_level)
    )

    total_ms = max(1, candidates["avg_monthly_searches"].sum())
    for _, row in candidates.iterrows():
        name = row["cluster_head"].strip().capitalize()
        ms = int(row["avg_monthly_searches"])
        score = round(ms / total_ms, 3)
        child = Node(
            id=gen_id(),
            parent_id=root.id,
            name=name,
            slug=f"/{make_slug(top_level_name)}/{make_slug(name)}",
            intent="transactional",
            ms=ms,
            score=score,
        )
        nodes.append(child)

    return nodes


def default_filters_for_node() -> List[FilterDef]:
    # Generic, safe defaults (can be edited later in Filters view)
    base = [
        ("Marca", "in"),
        ("Precio", "range"),
        ("Atributo 1", "in"),
        ("Atributo 2", "in"),
        ("Atributo 3", "in"),
    ]
    out: List[FilterDef] = []
    for i, (name, ftype) in enumerate(base, 1):
        out.append(FilterDef(name=name, type=ftype, values=[], order=i, visibility="core"))
    return out


def generate_metadata(node: Node, market: str) -> SEOData:
    # Template with guardrails
    titles = {
        "es": f"{node.name} al mejor precio | Envío rápido",
        "pt": f"{node.name} ao melhor preço | Envio rápido",
        "fr": f"{node.name} au meilleur prix | Expédition rapide",
        "it": f"{node.name} al miglior prezzo | Spedizione rapida",
    }
    metas = {
        "es": f"Compra {node.name} con envío rápido y garantía. Filtra por marca, precio y especificaciones.",
        "pt": f"Compre {node.name} com envio rápido e garantia. Filtre por marca, preço e especificações.",
        "fr": f"Achetez {node.name} avec expédition rapide et garantie. Filtrez par marque, prix et spécifications.",
        "it": f"Acquista {node.name} con spedizione rapida e garanzia. Filtra per marca, prezzo e specifiche.",
    }
    title = titles.get(market, titles["es"])[:65]
    meta = metas.get(market, metas["es"])[:160]
    h1 = node.name
    faq = [
        {"es": "¿Qué debo comparar antes de comprar?", "pt": "O que comparar antes de comprar?",
         "fr": "Que comparer avant d’acheter ?", "it": "Cosa confrontare prima di acquistare?"}[market if market in ["es","pt","fr","it"] else "es"],
        {"es": "¿Cómo filtrar por marca, precio y especificaciones?", "pt": "Como filtrar por marca, preço e especificações?",
         "fr": "Comment filtrer par marque, prix et spécifications ?", "it": "Come filtrare per marca, prezzo e specifiche?"}[market if market in ["es","pt","fr","it"] else "es"],
    ]
    return SEOData(title=title, h1=h1, meta=meta, faq=faq)


# -----------------------------
# Exporters
# -----------------------------
def export_master_json(nodes: List[Node]) -> str:
    obj = [asdict(n) for n in nodes]
    return json.dumps(obj, ensure_ascii=False, indent=2)


def export_taxonomy_csv(nodes: List[Node]) -> pd.DataFrame:
    return pd.DataFrame(
        [{"id": n.id, "parent_id": n.parent_id or "", "name": n.name, "slug": n.slug, "score": n.score} for n in nodes]
    )


def export_filters_csv(nodes: List[Node]) -> pd.DataFrame:
    rows = []
    for n in nodes:
        for f in n.filters:
            rows.append(
                {
                    "category_id": n.id,
                    "name": f.name,
                    "type": f.type,
                    "values": "|".join(f.values),
                    "order": f.order,
                    "visibility": f.visibility,
                }
            )
    return pd.DataFrame(rows)


def export_keywords_to_node(df_kw: pd.DataFrame, nodes: List[Node]) -> pd.DataFrame:
    """
    Map keywords to closest child node by simple TF-IDF cosine on cluster heads (POC).
    """
    if df_kw.empty or TfidfVectorizer is None:
        return pd.DataFrame(columns=["keyword", "node_id", "similarity", "intent"])

    # Build child names list
    root = next((n for n in nodes if n.parent_id is None), None)
    children = [n for n in nodes if n.parent_id == (root.id if root else None)]
    if not children:
        return pd.DataFrame(columns=["keyword", "node_id", "similarity", "intent"])

    # Vectorize
    texts = df_kw["kw_norm"].tolist() + [c.name for c in children]
    vec = TfidfVectorizer(ngram_range=(1, 2), min_df=1)
    X = vec.fit_transform(texts)
    kw_mat = X[: len(df_kw)]
    child_mat = X[len(df_kw):]

    sims = cosine_similarity(kw_mat, child_mat)
    best_idx = np.argmax(sims, axis=1)
    best_sim = sims[np.arange(len(df_kw)), best_idx]

    # Intent heuristic
    def intent_of(q: str) -> str:
        qn = normalize_text(q)
        if any(k in qn for k in ["precio", "oferta", "comprar", "barat", "mejor", "calidad precio",
                                  "preço", "oferta",
                                  "acheter", "prix", "promo",
                                  "prezzo", "comprare", "sconto"]):
            return "transactional"
        if any(k in qn for k in ["como", "cómo", "que es", "qué es", "cual", "cuál",
                                  "difference", "diferencia", "come", "cosa è", "o que é"]):
            return "informational"
        return "mixed"

    rows = []
    for i, kw in enumerate(df_kw["keyword"].tolist()):
        node = children[int(best_idx[i])]
        rows.append(
            {
                "keyword": kw,
                "node_id": node.id,
                "similarity": round(float(best_sim[i]), 3),
                "intent": intent_of(kw),
            }
        )
    return pd.DataFrame(rows)


def export_markdown(nodes: List[Node], gaps_df: pd.DataFrame) -> bytes:
    root = next((n for n in nodes if n.parent_id is None), None)
    children = [n for n in nodes if n.parent_id == (root.id if root else None)] if root else []
    md = io.StringIO()
    md.write(f"# Resumen ejecutivo — {root.name if root else 'Taxonomía'}\n\n")
    md.write("## Arquitectura propuesta\n```\n")
    if root:
        md.write(root.name + "\n")
        for c in children:
            md.write(f"└── {c.name}\n")
    else:
        md.write("(Árbol no disponible)\n")
    md.write("```\n\n")
    if not gaps_df.empty:
        md.write("## Top consultas (GSC) con impresiones\n\n")
        md.write(gaps_df.head(10).to_markdown(index=False))
    return md.getvalue().encode("utf-8")


def export_pptx(nodes: List[Node], coverage: Dict[str, int], gaps_df: pd.DataFrame) -> bytes:
    if not PPTX_AVAILABLE:
        return b""
    prs = Presentation()
    t = prs.slides.add_slide(prs.slide_layouts[0])
    t.shapes.title.text = "Resumen ejecutivo — Taxonomía"
    t.placeholders[1].text = "Arquitectura propuesta y oportunidades SEO/UX"

    # Coverage
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Cobertura por mercado (subset)"
    tf = s.placeholders[1].text_frame
    tf.text = ""
    if coverage:
        for mk, ms in coverage.items():
            p = tf.add_paragraph()
            p.text = f"• {mk.upper()}: {ms:,} búsquedas/mes"
    else:
        tf.text = "Sin datos de cobertura"

    # Tree
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Arquitectura propuesta"
    root = next((n for n in nodes if n.parent_id is None), None)
    children = [n for n in nodes if n.parent_id == (root.id if root else None)] if root else []
    tf2 = s2.placeholders[1].text_frame
    if root:
        tf2.text = root.name
        for c in children:
            p = tf2.add_paragraph()
            p.text = f"└── {c.name}"
    else:
        tf2.text = "(Árbol no disponible)"

    # Gaps (GSC)
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "GSC — Top consultas con impresiones"
    tf3 = s3.placeholders[1].text_frame
    tf3.text = ""
    if not gaps_df.empty:
        col = next((c for c in gaps_df.columns if "consulta" in c.lower() or "query" in c.lower()), gaps_df.columns[0])
        imp = next((c for c in ["impresiones", "impressions"] if c in gaps_df.columns), None)
        for _, r in gaps_df.head(8).iterrows():
            p = tf3.add_paragraph()
            text = str(r[col])
            if imp:
                text += f" — {int(r[imp])} impresiones"
            p.text = f"• {text}"
    else:
        tf3.text = "Sin datos de GSC"

    mem = io.BytesIO()
    prs.save(mem)
    mem.seek(0)
    return mem.read()


# -----------------------------
# Streamlit UI
# -----------------------------
def main():
    st.title("Generador de Taxonomías SEO (Tecnología)")
    st.caption("Taxonomía y facetado optimizados para SEO y UX · Multimercado ES/PT/FR/IT")

    with st.sidebar:
        st.subheader("Parámetros")
        preset_name = st.selectbox("Preset de calidad", list(PRESETS.keys()), index=0)
        preset = PRESETS[preset_name]
        max_levels = st.slider("Niveles máximos", 2, 3, preset["max_levels"])
        max_items = st.slider("Máx. ítems por nivel", 8, 12, preset["max_items_per_level"])
        st.text("")  # spacer

        st.subheader("Mercados")
        markets_enabled = st.multiselect("Selecciona mercados", SUPPORTED_MARKETS, default=["es", "pt", "fr", "it"])

        st.subheader("Nombres")
        top_level_by_market = {}
        for mk in markets_enabled:
            default_root = LOCALIZED_ROOT_DEFAULT.get(mk, "Tecnología")
            top_level_by_market[mk] = st.text_input(f"Nombre del nivel superior ({mk.upper()})", value=default_root, key=f"root_{mk}")

        st.divider()
        st.subheader("Guardar / Cargar sesión")
        load_json = st.file_uploader("Cargar sesión (.json)", type=["json"])
        if load_json is not None:
            try:
                st.session_state["session"] = json.load(load_json)
                st.success("Sesión cargada.")
            except Exception as e:
                st.error(f"No se pudo cargar la sesión: {e}")

    st.header("1) Carga de datos (GKP mínimo)")
    c1, c2, c3, c4 = st.columns(4)
    with c1:
        gkp_es = st.file_uploader("GKP España (CSV/Excel)", type=["csv", "xlsx"], key="gkp_es")
    with c2:
        gkp_pt = st.file_uploader("GKP Portugal (CSV/Excel)", type=["csv", "xlsx"], key="gkp_pt")
    with c3:
        gkp_fr = st.file_uploader("GKP Francia (CSV/Excel)", type=["csv", "xlsx"], key="gkp_fr")
    with c4:
        gkp_it = st.file_uploader("GKP Italia (CSV/Excel)", type=["csv", "xlsx"], key="gkp_it")

    st.markdown("**Opcionales**")
    oc1, oc2, oc3 = st.columns(3)
    with oc1:
        comp_paths = st.file_uploader("Estructuras/paths de competidores (CSV/Excel)", type=["csv", "xlsx"])
    with oc2:
        catalog = st.file_uploader("Catálogo interno (CSV/Excel)", type=["csv", "xlsx"])
    with oc3:
        gsc_zip = st.file_uploader("GSC export (ZIP)", type=["zip"])

    st.header("2) Ejecutar pipeline")
    run = st.button("Ejecutar / Re-agrupar")

    if run:
        with st.spinner("Procesando..."):
            # Build per-market dataframes
            market_files = {"es": gkp_es, "pt": gkp_pt, "fr": gkp_fr, "it": gkp_it}
            per_market: Dict[str, pd.DataFrame] = {}
            coverage: Dict[str, int] = {}

            for mk, f in market_files.items():
                if f is None or mk not in markets_enabled:
                    continue
                df = read_gkp_any(f)
                df = normalize_keywords(df)
                per_market[mk] = df
                coverage[mk] = int(df["avg_monthly_searches"].sum())

            if not per_market:
                st.warning("Carga al menos un GKP para continuar.")
                st.stop()

            # Optional: GSC gaps/cannibalization (ZIP export)
            gaps_df = pd.DataFrame()
            if gsc_zip is not None:
                qdf, _ = read_gsc_zip(gsc_zip)
                if not qdf.empty:
                    qcol = next((c for c in qdf.columns if "consulta" in c or "query" in c), qdf.columns[0])
                    imp = next((c for c in ["impresiones", "impressions"] if c in qdf.columns), None)
                    qdf["qnorm"] = qdf[qcol].astype(str).map(normalize_text)
                    # Show top by impressions to guide gap analysis (generic, not concept-bound)
                    gaps_df = qdf.sort_values(imp, ascending=False) if imp else qdf

            st.success("Datos cargados y normalizados por mercado.")

            # Build taxonomy per market
            all_nodes: Dict[str, List[Node]] = {}
            kw_map_all: Dict[str, pd.DataFrame] = {}

            for mk, df in per_market.items():
                # Clustering
                dfc = cluster_keywords(df, max_items=max_items)
                # Intent
                dfc = classify_intent(dfc)
                # Hierarchy (2-level for POC)
                top_level_name = top_level_by_market.get(mk, LOCALIZED_ROOT_DEFAULT.get(mk, "Tecnología"))
                nodes = propose_hierarchy_for_market(dfc, market=mk, top_level_name=top_level_name, max_items_per_level=max_items)
                # Filters + SEO for child nodes
                root = next((n for n in nodes if n.parent_id is None), None)
                for n in nodes:
                    if root and n.parent_id == root.id:
                        n.filters = default_filters_for_node()
                        n.seo = generate_metadata(n, market=mk)
                all_nodes[mk] = nodes

                # Keyword → node mapping (simple TF-IDF cosine to child names)
                kw_map = export_keywords_to_node(dfc[["keyword", "kw_norm"]], nodes)
                kw_map_all[mk] = kw_map

            st.header("3) Clustering & Taxonomía (gráfico + tabla)")
            for mk, nodes in all_nodes.items():
                st.subheader(f"Mercado: {mk.upper()}")
                df_nodes = pd.DataFrame(
                    [{"id": n.id, "parent_id": n.parent_id or "", "name": n.name, "slug": n.slug, "ms": n.ms, "score": n.score} for n in nodes]
                )
                if px is not None:
                    try:
                        fig = px.treemap(
                            df_nodes,
                            path=["parent_id", "name"],
                            values="ms",
                            hover_data=["score", "slug"],
                        )
                        st.plotly_chart(fig, use_container_width=True)
                    except Exception:
                        st.info("No se pudo renderizar el treemap; mostrando tabla.")
                        st.dataframe(df_nodes, use_container_width=True)
                else:
                    st.dataframe(df_nodes, use_container_width=True)

            st.header("4) Filtros y metadatos (resumen)")
            for mk, nodes in all_nodes.items():
                st.markdown(f"**{mk.upper()}**")
                child_rows = []
                for n in nodes:
                    if n.parent_id is not None:
                        child_rows.append(
                            {
                                "name": n.name,
                                "slug": n.slug,
                                "filters_core": ", ".join([f.name for f in n.filters if f.visibility == "core"]),
                                "title": n.seo.title,
                                "meta": n.seo.meta,
                            }
                        )
                if child_rows:
                    st.dataframe(pd.DataFrame(child_rows), use_container_width=True)

            st.header("5) Exportar")
            for mk, nodes in all_nodes.items():
                master_json = export_master_json(nodes)
                tax_csv = export_taxonomy_csv(nodes)
                filt_csv = export_filters_csv(nodes)
                kw_map = kw_map_all.get(mk, pd.DataFrame(columns=["keyword", "node_id", "similarity", "intent"]))

                st.download_button(I18N["download_json"]["es"].format(mk=mk), data=master_json.encode("utf-8"), mime="application/json", file_name=f"master_{mk}.json")
                st.download_button(I18N["download_tax"]["es"].format(mk=mk), data=tax_csv.to_csv(index=False).encode("utf-8"), mime="text/csv", file_name=f"taxonomy_{mk}.csv")
                st.download_button(I18N["download_filters"]["es"].format(mk=mk), data=filt_csv.to_csv(index=False).encode("utf-8"), mime="text/csv", file_name=f"filters_{mk}.csv")
                st.download_button(I18N["download_kw_map"]["es"].format(mk=mk), data=kw_map.to_csv(index=False).encode("utf-8"), mime="text/csv", file_name=f"keywords_to_node_{mk}.csv")

            md_bytes = export_markdown(next(iter(all_nodes.values())), gaps_df if not gaps_df.empty else pd.DataFrame())
            st.download_button(I18N["download_md"]["es"], data=md_bytes, mime="text/markdown", file_name="resumen_taxonomia.md")

            if PPTX_AVAILABLE:
                pptx_bytes = export_pptx(next(iter(all_nodes.values())), coverage, gaps_df if not gaps_df.empty else pd.DataFrame())
                st.download_button(I18N["download_pptx"]["es"], data=pptx_bytes, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", file_name="resumen_taxonomia.pptx")
            else:
                st.info("PPTX no disponible: instale `python-pptx` para habilitar la exportación a PowerPoint.")

            # Save session (lightweight)
            session_obj = {
                "preset": preset_name,
                "markets": list(per_market.keys()),
                "coverage": coverage,
                "timestamp": time.time(),
            }
            st.download_button("Guardar sesión (.json)", data=json.dumps(session_obj, ensure_ascii=False, indent=2).encode("utf-8"), mime="application/json", file_name="session_taxonomia.json")

    st.markdown("---")
    st.header("Comparativa con Google Search Console (opcional)")
    st.caption("Conecta tu propiedad (en próximas iteraciones) para detectar gaps y canibalización en tiempo real.")
    st.info("Stub listo: implementaremos OAuth y llamadas a la API cuando facilites las credenciales en `st.secrets`.")

    st.markdown("---")
    st.caption("MVP: clustering/intent/filtros simplificados. Próximas iteraciones: embeddings + HDBSCAN, árbitro LLM, edición inline, scraping SERP opcional.")


if __name__ == "__main__":
    main()
